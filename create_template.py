from pptx import Presentation
from pptx.util import Pt
import os

def create_template():
    # Create a new presentation
    prs = Presentation()

    # Add title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    # Set Pretendard font for title slide
    title.text = "Portfolio Template"
    for paragraph in title.text_frame.paragraphs:
        paragraph.font.name = 'Pretendard'
        paragraph.font.size = Pt(44)
        paragraph.font.bold = True

    subtitle.text = "Date Range"
    for paragraph in subtitle.text_frame.paragraphs:
        paragraph.font.name = 'Pretendard'
        paragraph.font.size = Pt(24)

    # Add content slide layout
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    
    # Set title
    title = slide.shapes.title
    title.text = "Project Template"
    for paragraph in title.text_frame.paragraphs:
        paragraph.font.name = 'Pretendard'
        paragraph.font.size = Pt(32)
        paragraph.font.bold = True

    # Add content placeholder
    body = slide.placeholders[1]
    tf = body.text_frame

    # Create directory if it doesn't exist
    os.makedirs('templates/pptx', exist_ok=True)
    
    # Save the template
    template_path = os.path.join('templates', 'pptx', 'test_template.pptx')
    prs.save(template_path)
    print(f"Template created at {template_path}")

if __name__ == "__main__":
    create_template()
